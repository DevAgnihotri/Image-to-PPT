from flask import Flask, render_template, request, send_file
from pptx import Presentation
from io import BytesIO

app = Flask(__name__)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload():
    files = request.files.getlist('file')
    if len(files) < 2:
        return "Please upload at least two images."
    
    prs = Presentation()
    slide_layout = prs.slide_layouts[6] # Two Content layout

    for i, file in enumerate(files):
        if i % 2 == 0:
            slide = prs.slides.add_slide(slide_layout)
            # Define the position relative to the top-left corner of the slide
            left_ = 50  # Horizontal position
            top_ = 50   # Vertical position
            left = slide.shapes.add_picture(BytesIO(file.read()), left_, top_, prs.slide_width / 2, prs.slide_height/1.2)
        else:
            right = slide.shapes.add_picture(BytesIO(file.read()), prs.slide_width / 1.8, 50, prs.slide_width / 2.3, prs.slide_height/1.2)



    ppt_bytes = BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)

    return send_file(ppt_bytes, as_attachment=True, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation', download_name='presentation.pptx')

if __name__ == '__main__':
    app.run(debug=True)
